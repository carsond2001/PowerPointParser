import re
from pptx import Presentation

# Open the PowerPoint file
prs = Presentation('pl12ch4.pptx')

# Open a text file to write the italicized words to
with open('italicized_words.txt', 'w', encoding='utf-8') as f:
    # Loop through each slide
    for slide in prs.slides:
        # Loop through each shape in the slide
        for shape in slide.shapes:
            # Check if the shape is a TextFrame
            if shape.has_text_frame:
                # Loop through each paragraph in the shape
                for paragraph in shape.text_frame.paragraphs:
                    # Loop through each run in the paragraph
                    for run in paragraph.runs:
                        # Check whether the run is italicized
                        if run.font.italic:
                            # Extract the text and write it to the file
                            italicized_text = run.text.strip()
                            # Replace non-ASCII characters with an appropriate string
                            italicized_text = re.sub(r'[^\x00-\x7F]+', ' ', italicized_text)
                            f.write(italicized_text + '\n')
